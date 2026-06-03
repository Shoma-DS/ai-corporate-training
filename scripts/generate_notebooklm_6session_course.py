#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import json
import re
import shutil
import sys
from pathlib import Path

sys.path.insert(0, "/private/tmp/ai_training_pydeps")

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
COURSE_NAME = "NotebookLMで進める業務効率化DXワークショップ講座"
COURSE_DIR = ROOT / "講座" / COURSE_NAME
OUT_DIR = ROOT / "書き出し" / "canva-pptx" / COURSE_NAME
TEMPLATE_ID = "soft-isometric-corporate-warm"
ACCESS_DATE = "2026-06-02"


SESSIONS = [
    {
        "no": "01",
        "dir": "01-業務課題整理とナレッジDX設計",
        "title": "業務課題整理とナレッジDX設計",
        "capability": "資料探索・確認業務を棚卸しし、NotebookLMに向く改善テーマを選定する",
        "output": "業務課題棚卸しシート、As-Is/To-Beメモ、活用テーマ選定メモ",
        "case": "返品問い合わせ対応の情報探索と初回回答準備",
    },
    {
        "no": "02",
        "dir": "02-情報源棚卸しと権限-機密区分設計",
        "title": "情報源棚卸しと権限・機密区分設計",
        "capability": "資料種別、更新日、責任者、機密区分、取り込み可否を判断する",
        "output": "情報源棚卸し表、入力禁止情報チェック、権限・保存設計メモ",
        "case": "FAQ、マニュアル、会議メモ、公開資料の取り込み可否判断",
    },
    {
        "no": "03",
        "dir": "03-根拠付き質問とFAQ-マニュアル整備",
        "title": "根拠付き質問とFAQ・マニュアル整備",
        "capability": "根拠付き回答を使い、FAQ、確認フロー、マニュアル改善案を作成する",
        "output": "根拠付きFAQ、社内確認チェックリスト、マニュアル改善メモ",
        "case": "返品条件と例外対応を顧客向けFAQと社内確認フローに分ける",
    },
    {
        "no": "04",
        "dir": "04-会議-研修-営業資料の要約とアクション化",
        "title": "会議・研修・営業資料の要約とアクション化",
        "capability": "議事録、音声、動画、スライド資料を要約し、担当・期限・次アクションへ変換する",
        "output": "会議アクション一覧、営業説明文案、研修理解メモ",
        "case": "営業会議メモと製品説明資料から次アクションと説明文案を作る",
    },
    {
        "no": "05",
        "dir": "05-共有-レビュー-継続運用設計",
        "title": "共有・レビュー・継続運用設計",
        "capability": "共有範囲、更新頻度、レビュー担当、誤回答対応、削除ルールを設計する",
        "output": "ナレッジ運用設計書、リスクチェックリスト、レビュー手順書",
        "case": "部署内でNotebookLMを安全に使い続ける運用ルール作成",
    },
    {
        "no": "06",
        "dir": "06-業務効率化DX提案書と展開計画",
        "title": "業務効率化DX提案書と展開計画",
        "capability": "削減時間、品質安定、定着施策、導入ロードマップを提案書にまとめる",
        "output": "NotebookLM業務効率化DX提案書、導入ロードマップ、KPI試算表",
        "case": "小さな導入テーマを部署展開の提案書にする",
    },
]


BASE_SLIDES = [
    ("{title}", "第{no}回の到達点と成果物を確認する", "summary-takeaway"),
    ("この回で作る成果物", "{output}", "checklist-review"),
    ("業務課題から始める理由", "機能紹介ではなく、確認・説明・共有の業務改善として扱う", "before-after-cards"),
    ("6回ロードマップ上の位置づけ", "前回成果物との接続と次回への受け渡しを確認する", "soft-process-flow"),
    ("情報管理の前提", "実データ、個人情報、契約情報、未公開資料を演習に使わない", "caution-risk-check"),
    ("今回の業務ケース", "{case}", "role-card-set"),
    ("Before: 現場で起きている詰まり", "資料分散、最新版不明、根拠確認漏れ、担当者依存", "before-after-cards"),
    ("After: NotebookLMを業務に組み込む姿", "情報源、質問、確認、共有、更新を一連の流れにする", "soft-process-flow"),
    ("画面共有: 講師記入例を見る", "この回で使うサンプル資料とワークシートを確認する", "screen-share-transition"),
    ("作業の全体像", "入力、処理、出力、確認、保存、共有、改善の流れ", "soft-process-flow"),
    ("判断基準1: NotebookLMに向く業務", "複数資料を横断し、根拠付きで確認したい業務", "permission-role-matrix"),
    ("判断基準2: 任せない業務", "承認、例外判断、社外送信、機密情報を含む判断", "caution-risk-check"),
    ("ワーク1: 現状を書き出す", "動画を一時停止して、課題と資料をワークシートへ記入する", "workshop-instruction"),
    ("講師記入例の確認", "記入粒度、抽象化、公開不可情報の除外を確認する", "data-table-checkpoints"),
    ("自己レビュー1", "業務課題、資料、困りごと、成果物がつながっているか", "checklist-review"),
    ("資料を情報源として設計する", "資料名、形式、目的、更新日、責任者、機密区分を揃える", "data-table-checkpoints"),
    ("質問は目的から逆算する", "読み手、出力形式、根拠確認方法をセットで設計する", "soft-process-flow"),
    ("根拠付き回答の見方", "引用、根拠資料、未確認事項、追加確認先を分ける", "checklist-review"),
    ("画面共有: NotebookLMで質問する", "ダミー資料を使って根拠付き回答の作り方を見る", "screen-share-transition"),
    ("良い質問と悪い質問", "曖昧な依頼を、検証可能な質問へ直す", "ui-mockup-comparison"),
    ("出力を業務文書へ変換する", "FAQ、説明文、アクション一覧、提案メモに整える", "soft-process-flow"),
    ("人が確認する位置を残す", "AI下書き、原典確認、担当者レビュー、最終判断を分ける", "caution-risk-check"),
    ("ワーク2: 根拠付き出力を作る", "動画を一時停止して、指定の成果物を作成する", "workshop-instruction"),
    ("講師記入例との比較", "根拠、未確認、追加確認が分かれているか確認する", "data-table-checkpoints"),
    ("よくある失敗", "古い資料、過剰な断定、根拠なし回答、共有範囲ミス", "caution-risk-check"),
    ("Studio出力の使い分け", "ブリーフィング、Mind Map、Audio Overview等を目的別に扱う", "role-card-set"),
    ("共有前チェック", "誰に見せてよいか、何を削るか、誰が承認するか", "collaboration-scope-check"),
    ("更新ルールを決める", "更新担当、更新頻度、差し替え、削除、停止条件", "folder-structure"),
    ("画面共有: 運用ルールの記入例", "運用設計シートを講師例で確認する", "screen-share-transition"),
    ("リスクを一覧化する", "情報漏洩、著作権、ハルシネーション、権限過多", "caution-risk-check"),
    ("効果を説明できる形にする", "探索時間、一次回答作成、会議後整理、品質安定", "data-table-checkpoints"),
    ("小さく始める導入ステップ", "1業務、1部署、1ノートブック、1か月で試す", "soft-process-flow"),
    ("ワーク3: 活用設計を完成する", "動画を一時停止して、対象業務と運用ルールを記入する", "workshop-instruction"),
    ("講師記入例の最終確認", "成果物の抜け漏れと安全性を確認する", "checklist-review"),
    ("自部署へ置き換える観点", "資料、役割、承認、更新、KPIを自部署名なしで整理する", "role-card-set"),
    ("ワーク4: 次アクションを決める", "最初に試すテーマ、確認先、期限、レビュー方法を書く", "workshop-instruction"),
    ("第{no}回の成果物チェック", "{output}", "checklist-review"),
    ("次回への受け渡し", "今回の成果物を次回の設計・実演・提案へ接続する", "soft-process-flow"),
    ("受講後の自己点検", "公開不可情報、根拠確認、人のレビュー、更新ルール", "caution-risk-check"),
    ("まとめ", "{capability}", "summary-takeaway"),
]


SESSION_EXTRA = {
    "01": {
        10: "業務を入力、処理、出力、確認、保存、共有、改善に分解する",
        16: "As-Is/To-Beを資料探索業務に適用する",
        23: "自部署の情報探索業務を1つ選んでAs-Is/To-Beを作る",
        37: "業務棚卸し、As-Is/To-Be、活用テーマを確認する",
    },
    "02": {
        10: "情報源を資料名、形式、更新日、責任者、機密区分で管理する",
        16: "Drive、PDF、Docs、Slides、Web、YouTube、音声の扱いを分ける",
        23: "資料一覧CSVから取り込み可否と事前処理を判断する",
        37: "情報源棚卸し表と入力禁止情報チェックを確認する",
    },
    "03": {
        10: "質問設計を目的、読み手、出力形式、根拠確認で組み立てる",
        16: "返品FAQ素材から顧客向けFAQと社内確認フローを分ける",
        23: "根拠付きFAQを作り、断定しすぎた表現を直す",
        37: "根拠付きFAQ、社内確認チェック、改善メモを確認する",
    },
    "04": {
        10: "会議メモ、研修資料、営業資料をアクション化する手順を確認する",
        16: "担当、期限、決定事項、未決事項、確認先を抽出する",
        23: "会議メモからアクション一覧と営業説明文案を作る",
        37: "会議アクション一覧、説明文案、理解メモを確認する",
    },
    "05": {
        10: "共有範囲、更新担当、レビュー担当、削除条件を設計する",
        16: "Workspace利用時の管理者設定、コピー、DLP、CAA、Takeoutを確認する",
        23: "ナレッジ運用設計書とリスクチェックリストを作る",
        37: "運用設計書、レビュー手順、リスクチェックを確認する",
    },
    "06": {
        10: "削減時間、品質安定、属人化解消を提案書の言葉にする",
        16: "KPI効果試算と導入ロードマップを組み立てる",
        23: "NotebookLM業務効率化DX提案書の骨子を作る",
        37: "提案書、ロードマップ、KPI試算表を確認する",
    },
}


FONT_CACHE: dict[tuple[int, bool], ImageFont.ImageFont] = {}


def font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    key = (size, bold)
    if key in FONT_CACHE:
        return FONT_CACHE[key]
    candidates = [
        "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc" if bold else "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    for path in candidates:
        try:
            FONT_CACHE[key] = ImageFont.truetype(path, size=size)
            return FONT_CACHE[key]
        except OSError:
            continue
    FONT_CACHE[key] = ImageFont.load_default()
    return FONT_CACHE[key]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.ImageFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        current = ""
        for ch in paragraph:
            trial = current + ch
            if draw.textbbox((0, 0), trial, font=fnt)[2] <= max_width:
                current = trial
            else:
                if current:
                    lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def draw_card(draw: ImageDraw.ImageDraw, xy: tuple[int, int, int, int], fill: str, outline: str = "#D8E3EA") -> None:
    draw.rounded_rectangle(xy, radius=24, fill=fill, outline=outline, width=2)


def render_slide(path: Path, session: dict[str, str], slide_no: int, title: str, body: str, pattern: str) -> None:
    img = Image.new("RGB", (1920, 1080), "#FBFCF8")
    d = ImageDraw.Draw(img)
    title_f = font(58, True)
    sub_f = font(30)
    small_f = font(24)
    body_f = font(34)
    d.rectangle((0, 0, 1920, 1080), fill="#FBFCF8")
    d.rounded_rectangle((64, 48, 1860, 152), radius=28, fill="#FFFFFF", outline="#DDE7EF", width=2)
    d.text((92, 72), f"S{slide_no:02d}  {title}", fill="#16324F", font=title_f)
    d.text((92, 158), f"{COURSE_NAME} / 第{session['no']}回 {session['title']}", fill="#4B6B84", font=small_f)
    d.rounded_rectangle((1540, 64, 1828, 128), radius=20, fill="#E7F3F1", outline="#BBDDD8", width=2)
    d.text((1574, 78), "2時間 / 録画ワーク", fill="#0F6B65", font=small_f)

    colors = ["#E8F2FF", "#EAF7F0", "#FFF3DC", "#F1ECFF", "#FDECEC"]
    labels = [
        "業務課題",
        "情報源",
        "根拠確認",
        "人のレビュー",
        "運用改善",
    ]
    for i, label in enumerate(labels):
        x = 116 + i * 350
        draw_card(d, (x, 250, x + 270, 420), colors[i])
        d.ellipse((x + 24, 282, x + 76, 334), fill=["#2F80ED", "#1B9A70", "#E39A21", "#7B61D1", "#D85858"][i])
        d.text((x + 96, 280), label, fill="#16324F", font=body_f)
        if i < len(labels) - 1:
            d.line((x + 284, 335, x + 330, 335), fill="#7AA9C4", width=6)
            d.polygon([(x + 330, 335), (x + 314, 324), (x + 314, 346)], fill="#7AA9C4")

    draw_card(d, (116, 500, 1240, 900), "#FFFFFF")
    d.text((156, 534), "このスライドの要点", fill="#16324F", font=body_f)
    wrapped = wrap_text(d, body, sub_f, 1020)
    y = 604
    for line in wrapped[:7]:
        d.text((170, y), "・" + line, fill="#243B53", font=sub_f)
        y += 48
    draw_card(d, (1300, 500, 1800, 900), "#F4F8FB")
    d.text((1340, 534), "図解パターン", fill="#16324F", font=body_f)
    d.text((1340, 604), pattern, fill="#2F80ED", font=sub_f)
    d.text((1340, 690), "Canvaではこの画像を\nページ全面に配置し、\n表紙・章扉などのみ\n必要に応じて編集化。", fill="#4B5563", font=sub_f)
    d.text((92, 1018), "Public-safe dummy course material. No personal data, no customer data, no credentials.", fill="#8796A5", font=small_f)
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, compress_level=1)


def slide_rows(session: dict[str, str]) -> list[dict[str, str | int]]:
    rows = []
    extras = SESSION_EXTRA[session["no"]]
    for idx, (title, body, pattern) in enumerate(BASE_SLIDES, start=1):
        if idx in extras:
            body = extras[idx]
        rows.append(
            {
                "no": idx,
                "title": title.format(**session),
                "body": body.format(**session),
                "pattern": pattern,
            }
        )
    return rows


def write(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def write_csv(path: Path, rows: list[list[str]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerows(rows)


def session_slide_plan(session: dict[str, str], rows: list[dict[str, str | int]]) -> str:
    lines = [
        f"# 第{int(session['no'])}回 スライド構成案: {session['title']}",
        "",
        f"対象講座: {COURSE_NAME}",
        "",
        "## 設計方針",
        "",
        "- 120分の録画eラーニングとして、講義、画面共有、個人ワーク、講師記入例との自己レビューを組み合わせる。",
        f"- 採用テンプレート: `{TEMPLATE_ID}`",
        f"- この回の業務遂行力: {session['capability']}",
        f"- 成果物: {session['output']}",
        "- 実在の社名、顧客名、個人情報、価格、連絡先、社内未公開資料は使わない。",
        "",
        "## 120分の時間配分",
        "",
        "| ブロック | 時間 | スライド範囲 | ねらい |",
        "| --- | ---: | --- | --- |",
        "| 導入と全体像 | 15分 | S01-S05 | 到達点、成果物、情報管理の前提を確認する |",
        "| 業務ケースと判断基準 | 25分 | S06-S15 | 業務課題、適用範囲、ワーク1を行う |",
        "| 情報源と質問設計 | 30分 | S16-S25 | 資料設計、根拠付き質問、ワーク2を行う |",
        "| 運用・効果・設計完成 | 35分 | S26-S36 | 共有、更新、リスク、KPI、ワーク3-4を行う |",
        "| まとめと次回接続 | 15分 | S37-S40 | 成果物確認と次回準備 |",
        "",
        "## スライド一覧",
        "",
        "| No. | タイトル | 主な内容 | 使用テンプレート | 図解パターン | 素材 |",
        "| ---: | --- | --- | --- | --- | --- |",
    ]
    for row in rows:
        material = "画面共有案内" if "画面共有" in str(row["title"]) else "図解/ワーク"
        lines.append(f"| S{row['no']:02d} | {row['title']} | {row['body']} | {TEMPLATE_ID} | {row['pattern']} | {material} |")
    lines += [
        "",
        "## 画面共有で見せる箇所",
        "",
        "| タイミング | 見せるもの | 目的 |",
        "| --- | --- | --- |",
        "| 0:25頃 | `演習データ/README.md` と講師記入例 | 演習で使うダミー資料の扱いを確認する |",
        "| 0:55頃 | NotebookLMのダミーノートブックまたは代替資料 | 根拠付き質問と引用確認の流れを示す |",
        "| 1:25頃 | ワークシート記入例 | 出力確認と運用設計の粒度を示す |",
        "",
    ]
    return "\n".join(lines)


def session_script(session: dict[str, str], rows: list[dict[str, str | int]]) -> str:
    lines = [
        f"# 第{int(session['no'])}回 講師台本: {session['title']}",
        "",
        f"講座名: {COURSE_NAME}",
        "",
        "## この台本の使い方",
        "",
        "この台本は録画動画によるeラーニング用です。ライブ発表、チャット共有、相互レビューは行いません。ワークは動画を一時停止して取り組む形で進めます。",
        "",
        "## 準備物",
        "",
        "- 第{no}回スライド S01-S40".format(**session),
        "- `ワークシート.md`",
        "- `配布資料/演習ガイド.md`",
        "- `演習データ/README.md`",
        "",
        "講師メモ:",
        f"第{int(session['no'])}回は「{session['capability']}」を中心に進める。NotebookLMの機能紹介に寄りすぎた場合は、業務課題、根拠確認、人のレビュー、運用改善へ戻す。",
        "",
    ]
    for row in rows:
        n = int(row["no"])
        title = str(row["title"])
        body = str(row["body"])
        lines += [
            "スライド切替:",
            f"S{n:02d}「{title}」",
            "",
        ]
        if n in [9, 19, 29]:
            lines += [
                f"画面共有 ── 実演{[9,19,29].index(n)+1}「{title.replace('画面共有: ', '')}」",
                "⏱ 約5分",
                "",
                "【手順1 – 約1分】",
                f"`演習データ/README.md` とこの回のワークシートを開き、{session['case']} に関係するダミー資料だけを使うことを説明する。",
                "",
                "【手順2 – 約1分】",
                "NotebookLMを利用できる環境では、ダミー資料をノートブックに追加する。利用できない環境では、配布資料の本文を読みながら質問設計だけを行う代替手順を示す。",
                "",
                "【手順3 – 約1分】",
                "質問文、期待する出力、根拠確認方法を声に出して確認する。回答をそのまま使わず、原典確認と人のレビューに分ける。",
                "",
                "【手順4 – 約1分】",
                "講師記入例を見せ、根拠、未確認事項、追加確認先が分かれているか確認する。",
                "",
                "【手順5 – 約1分】",
                "この実演が、次のワークで作る成果物につながることを説明する。",
                "",
                "【見せるポイント】",
                "NotebookLMの操作そのものではなく、資料の選定、質問の設計、出力確認、運用ルールまでを一つの業務フローとして見せる。",
                "",
            ]
        elif n in [13, 23, 33, 36]:
            minutes = "10" if n in [13, 23] else "8"
            lines += [
                "ワーク指示:",
                f"「ここで動画を一時停止して、{minutes}分ほど取り組んでください。{body}。実在の顧客名、個人名、メールアドレス、価格、契約情報は書かず、ダミーまたは抽象化した表現で記入します。取り組めたら再生してください。」",
                "",
            ]
        elif n in [5, 12, 22, 25, 30, 39]:
            lines += [
                "読み上げ:",
                f"「{body}。ここは情報管理と業務品質に関わる重要な確認です。NotebookLMは資料に基づく回答を支援しますが、最終判断は人が行います。根拠、未確認事項、追加確認先を分けて残してください。」",
                "",
            ]
        else:
            lines += [
                "読み上げ:",
                f"「{body}。この回では、{session['capability']}ことを目指します。最終的には、{session['output']}として持ち帰れる形に整えます。」",
                "",
            ]
    return "\n".join(lines)


def prompt_file(session: dict[str, str], rows: list[dict[str, str | int]]) -> str:
    lines = [
        f"# 第{int(session['no'])}回 画像生成プロンプト",
        "",
        f"- 採用テンプレート: `{TEMPLATE_ID}`",
        "- 画角: 16:9、1920x1080px",
        "- 注意: 実在ロゴ、実在UI、個人情報、価格、連絡先、APIキー、社内資料原文は画像生成AIに描かせない。",
        "",
    ]
    for row in rows:
        lines += [
            f"## S{row['no']:02d}. {row['title']}",
            "",
            f"- 使用テンプレート: `{TEMPLATE_ID}`",
            f"- 使用図解パターン: `{row['pattern']}`",
            "- 使用素材: 図解。実演スライドは画面共有案内のみ。",
            f"- 画像内テキスト: `{row['title']}`",
            "- 画像プロンプト:",
            "",
            "```text",
            "Create a 16:9 Japanese corporate training slide in a soft warm isometric vector illustration style. Use a clean white or slightly warm off-white background, large bold Japanese title area, gentle rounded cards, thin light-gray borders, soft shadows, friendly simplified people illustrations, clear tables or diagrams, navy headline text, blue primary accents, green success accents, red caution accents, and small orange or purple highlights only when useful.",
            f"Slide title: {row['title']}",
            f"Key message: {row['body']}",
            f"Diagram pattern: {row['pattern']}",
            "No real brand logos, no real UI screenshots, no confidential information, no personal data, no pricing, no contact details, no placeholder boxes.",
            "```",
            "",
        ]
    return "\n".join(lines)


def worksheet(session: dict[str, str]) -> str:
    return f"""# 第{int(session['no'])}回 ワークシート: {session['title']}

## 1. この回の対象業務

| 項目 | 記入欄 |
| --- | --- |
| 対象業務 |  |
| 使う資料 |  |
| 作りたい成果物 | {session['output']} |
| 情報管理上の注意 |  |

## 2. 資料・根拠確認

| 資料名 | 形式 | 更新日 | 責任者 | 機密区分 | 取り込み可否 | 根拠確認方法 |
| --- | --- | --- | --- | --- | --- | --- |
|  |  |  |  |  |  |  |

## 3. 質問設計

| 目的 | 質問文 | 期待する出力 | 確認観点 |
| --- | --- | --- | --- |
|  |  |  |  |

## 4. 作成した成果物

記入:

## 5. 自己レビュー

| 観点 | はい/いいえ | 修正メモ |
| --- | --- | --- |
| 根拠資料を確認した |  |  |
| 未確認事項を分けた |  |  |
| 個人情報・顧客情報を書いていない |  |  |
| 人のレビュー位置を決めた |  |  |
| 更新担当・共有範囲を決めた |  |  |
"""


def exercise_guide(session: dict[str, str]) -> str:
    return f"""# 第{int(session['no'])}回 演習ガイド

## この回の成果物

{session['output']}

## 演習ケース

{session['case']}

## 進め方

1. `演習データ/README.md` を読み、使うダミー資料を確認する。
2. NotebookLMを利用できる場合は、ダミー資料をノートブックに追加する。
3. 利用できない場合は、配布資料を読みながら質問文と期待出力を設計する。
4. ワークシートに根拠、未確認事項、追加確認先を分けて記入する。
5. 講師記入例と比較し、公開不可情報が含まれていないか確認する。

## 注意

実在の顧客情報、社員情報、契約情報、価格、連絡先、未公開資料、APIキーは使わない。
"""


def exercise_readme(session: dict[str, str]) -> str:
    return f"""# 第{int(session['no'])}回 演習データ

このフォルダのデータは、{COURSE_NAME} 第{int(session['no'])}回用の公開可能なダミーデータです。

## ケース

{session['case']}

## 使うファイル

- `資料棚卸しサンプル.csv`
- `業務メモサンプル.md`
- `講師記入例.md`

## 注意

実在企業、実在顧客、実在社員の情報は含めていません。実務適用時は自社ルール、管理者設定、共有範囲、入力禁止情報を確認してください。
"""


def create_pptx(session_dir: Path, session: dict[str, str], rows: list[dict[str, str | int]]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for row in rows:
        slide = prs.slides.add_slide(blank)
        img_path = session_dir / "スライド画像" / f"S{int(row['no']):02d}.png"
        slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"{session['no']}-{session['title']}.pptx"
    prs.save(path)
    return path


def write_course_files() -> None:
    whole_curriculum = "\n".join(
        [
            "| 回 | テーマ | 時間 | 主な内容 | 成果物 |",
            "| ---: | --- | ---: | --- | --- |",
            *[
                f"| {int(s['no'])} | {s['title']} | 120分 | {s['capability']} | {s['output']} |"
                for s in SESSIONS
            ],
        ]
    )
    write(
        COURSE_DIR / "全体" / "講座概要.md",
        f"""# 講座概要

## タイトルと一言訴求

{COURSE_NAME}。

散在する社内資料、FAQ、マニュアル、議事録、提案資料を、根拠付きで読み解き、現場で使える業務アウトプットと継続運用ルールへ変える12時間のリスキリング講座。

## 研修の目的

NotebookLMを単なる要約ツールではなく、業務課題整理、情報源設計、根拠付き質問、出力確認、共有範囲判断、更新運用、DX提案までを支える業務プロセスとして扱う。

## 講座区分

- 想定区分: リスキリングコース
- 標準学習時間: 6回、各120分、合計約12時間
- 標準受講期間: 2か月以内
- 提供方式: 録画動画によるeラーニング
- 受講管理: LMSによる視聴状況、受講時間、課題提出、修了確認の記録

## 6回の構成

{whole_curriculum}

## 受講後にできること

- NotebookLMに向く業務と向かない業務を判断できる。
- 情報源の更新日、責任者、機密区分、取り込み可否を整理できる。
- 根拠付きFAQ、会議アクション一覧、営業説明文案、運用設計書を作れる。
- AI出力を原典確認、人のレビュー、未確認事項、追加確認先に分けられる。
- 自部署で継続運用するための更新ルール、共有範囲、KPI、導入ロードマップを提案できる。

## 注意事項

NotebookLMの機能、利用可否、上限、共有設定はアカウント種別、契約、管理者設定、国・地域により変わる。実務適用時は最新の公式情報と自社の情報管理ルールを確認する。
""",
    )
    write(
        COURSE_DIR / "全体" / "詳細シラバス.md",
        f"""# 詳細シラバス

## 到達目標

受講者は、NotebookLMを使って業務資料を根拠付きで読み解き、FAQ、説明文、会議アクション、ナレッジ運用設計、DX提案書を作成できる。

## 12時間カリキュラム

{whole_curriculum}

## 評価観点

- 業務課題からNotebookLM活用テーマを選定できているか。
- 資料の機密区分、更新日、責任者、共有範囲を確認しているか。
- NotebookLMの回答を根拠、引用、未確認事項、追加確認先に分けているか。
- AI出力をそのまま使わず、人による確認、原典確認、例外対応を組み込んでいるか。
- 自部署で継続運用するための更新ルール、レビュー担当、削除ルールがあるか。
""",
    )
    write(
        COURSE_DIR / "全体" / "演習データ回別一覧.md",
        "# 演習データ回別一覧\n\n"
        + "\n".join(
            [
                "| 回 | フォルダ | 主な演習データ | 用途 |",
                "| ---: | --- | --- | --- |",
                *[
                    f"| {int(s['no'])} | `{s['dir']}/演習データ/` | 資料棚卸しサンプル.csv、業務メモサンプル.md、講師記入例.md | {s['case']} |"
                    for s in SESSIONS
                ],
            ]
        )
        + "\n",
    )
    write(
        COURSE_DIR / "全体" / "スライド構成案.md",
        "# 全体スライド構成案\n\n"
        f"- 採用テンプレート: `{TEMPLATE_ID}`\n"
        "- 各回40枚、全240枚\n"
        "- Canva取り込みは `書き出し/canva-pptx/` の回別PPTXを使用する。\n\n"
        + whole_curriculum
        + "\n",
    )
    write(
        COURSE_DIR / "全体" / "Canva納品メモ.md",
        f"""# Canva納品メモ

## 方針

`corporate-training-course-builder` と `gws-ai-training-slide-exporter` の方針に従い、画像アップロード先行のCanva納品を標準にする。

## ローカル成果物

- 回別PPTX: `書き出し/canva-pptx/{COURSE_NAME}/`
- 回別スライド画像: 各回 `スライド画像/S01.png` から `S40.png`

## Canva上の推奨構成

Canvaフォルダ:

```text
AI法人研修/
  {COURSE_NAME}/
    01-業務課題整理とナレッジDX設計
    02-情報源棚卸しと権限-機密区分設計
    03-根拠付き質問とFAQ-マニュアル整備
    04-会議-研修-営業資料の要約とアクション化
    05-共有-レビュー-継続運用設計
    06-業務効率化DX提案書と展開計画
```

## Magic Layers対象

全ページではなく、表紙、章扉、成果物一覧、カリキュラム表、提案書ページのみ手動Magic Layers化する。

## URL管理

Canva URL、design ID、編集メモは `非公開/Canva/` に保存し、publicリポジトリには書かない。
""",
    )


def create_session(session: dict[str, str], *, skip_pptx: bool = False, skip_existing_images: bool = False) -> Path:
    session_dir = COURSE_DIR / session["dir"]
    for sub in ["スライド画像", "スクリーンショット", "配布資料", "演習データ"]:
        (session_dir / sub).mkdir(parents=True, exist_ok=True)
    rows = slide_rows(session)
    write(session_dir / "スライド案.md", session_slide_plan(session, rows))
    write(session_dir / "講師台本.md", session_script(session, rows))
    write(session_dir / "画像生成プロンプト.md", prompt_file(session, rows))
    write(session_dir / "ワークシート.md", worksheet(session))
    write(session_dir / "配布資料" / "演習ガイド.md", exercise_guide(session))
    write(session_dir / "演習データ" / "README.md", exercise_readme(session))
    write(
        session_dir / "演習データ" / "業務メモサンプル.md",
        f"# 業務メモサンプル\n\nケース: {session['case']}\n\nこの資料は架空企業のダミー資料です。顧客名、個人名、価格、連絡先は含めません。\n\n## 状況\n\n資料が複数に分散し、担当者が根拠を探すのに時間がかかっています。NotebookLMを使い、根拠付きの確認、説明文案、次アクション整理へつなげます。\n\n## 未確認事項\n\n- 最新版資料の確認\n- 共有範囲の確認\n- 担当者レビューの位置\n",
    )
    write(
        session_dir / "演習データ" / "講師記入例.md",
        f"# 講師記入例\n\n## 成果物\n\n{session['output']}\n\n## 記入例\n\n- 対象業務: {session['case']}\n- 入れる資料: 公開可能なFAQ、匿名化した会議メモ、更新日が確認できるマニュアル\n- 入れない資料: 個人情報、取引条件、価格表、未公開資料\n- 確認方法: 根拠資料、未確認事項、追加確認先を分けて記録する\n",
    )
    write_csv(
        session_dir / "演習データ" / "資料棚卸しサンプル.csv",
        [
            ["資料名", "形式", "用途", "更新日", "責任者", "機密区分", "取り込み可否", "事前処理"],
            ["公開FAQ", "Docs", session["case"], "2026-05-15", "業務管理", "公開可", "可", "最新版確認"],
            ["社内マニュアル抜粋", "PDF", "確認フロー", "2026-04-01", "運用担当", "社内限定", "条件付き", "機密箇所を抜粋"],
            ["会議メモ要約版", "Markdown", "アクション整理", "2026-05-25", "プロジェクト担当", "社内限定", "条件付き", "個人名を役割名へ置換"],
            ["取引条件表", "Sheet", "個別判断", "2026-05-10", "営業管理", "機密", "不可", "使用しない"],
        ],
    )
    for row in rows:
        image_path = session_dir / "スライド画像" / f"S{int(row['no']):02d}.png"
        if skip_existing_images and image_path.exists():
            continue
        print(f"render {session['no']} S{int(row['no']):02d}", flush=True)
        render_slide(image_path, session, int(row["no"]), str(row["title"]), str(row["body"]), str(row["pattern"]))
    if not skip_pptx:
        create_pptx(session_dir, session, rows)
    return session_dir


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--skip-pptx", action="store_true")
    parser.add_argument("--pptx-only", action="store_true")
    parser.add_argument("--skip-existing-images", action="store_true")
    args = parser.parse_args()
    if args.pptx_only:
        for session in SESSIONS:
            create_pptx(COURSE_DIR / session["dir"], session, slide_rows(session))
        print(f"pptx output: {OUT_DIR}")
        return
    write_course_files()
    created = [create_session(session, skip_pptx=args.skip_pptx, skip_existing_images=args.skip_existing_images) for session in SESSIONS]
    write(
        OUT_DIR / "manifest.json",
        json.dumps(
            {
                "course": COURSE_NAME,
                "template_id": TEMPLATE_ID,
                "sessions": [
                    {
                        "no": s["no"],
                        "title": s["title"],
                        "pptx": str((OUT_DIR / f"{s['no']}-{s['title']}.pptx").relative_to(ROOT)),
                        "slides": 40,
                    }
                    for s in SESSIONS
                ],
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
    )
    print(f"created {len(created)} sessions under {COURSE_DIR}")
    print(f"pptx output: {OUT_DIR}")


if __name__ == "__main__":
    main()
